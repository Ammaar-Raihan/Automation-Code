# This code requests a graph from the URLs mentioned in image_dictionary, edits the images as required and replaces them in the report powerpoint

# 21/1/2024
# The code has been edited to remove any mentions of production file paths and URLs

import time
from pptx import Presentation
import datetime
import requests
from PIL import Image

image_dictionary = {
    "Slide 3 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 3 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 4 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 4 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 5 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 5 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 6 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 6 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 7 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 7 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 8 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 8 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 9 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 9 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 10 CPU": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
    "Slide 10 Mem": "prtg.url/chart.png?graphid=-1&id=ID&avg=86400&sdate=STARTDATE-00-00-00&edate=ENDDATE-23-59-00&clgid=&width=850&height=270&deftheme=1&graphstyling=baseFontSize=%2711%27%20&refreshable=true",
}

# create a requests sessions
session = requests.Session()

try:
    with open(r'.\key.txt', 'r') as file:
        text = file.read()
except FileNotFoundError:
    print("File not found.")
except IOError:
    print("Error reading the file.")

API_TOKEN = text
file.close()
# Prompt the user for a month number
month = int(input("Enter a month number (1-12): "))

# Get the current year and month
current_year = 2023
current_month = datetime.datetime.now().month

# # Check if the input is greater than the previous month
# if month > current_month-1:
#     print("Input cannot be greater than the previous month.")
#     exit()

# Create a datetime object for the first day of the month
first_day = datetime.datetime(current_year, month, 1)

# Find the last day of the month by adding one month and subtracting one day
next_month = first_day.replace(day=28) + datetime.timedelta(days=4)
last_day = next_month - datetime.timedelta(days=next_month.day)

# Format the dates as YYYY-MM-DD
first_day_str = first_day.strftime("%Y-%m-%d")
last_day_str = last_day.strftime("%Y-%m-%d")

# Print the dates
print("First day of the month:", first_day_str)
print("Last day of the month:", last_day_str)

for key, value in image_dictionary.items():
    value = value.replace("STARTDATE", first_day_str)
    value = value.replace("ENDDATE", last_day_str)
    image_dictionary[key] = value

# Load the PowerPoint presentation
presentation = Presentation(r"file\path\report.pptx")

# Iterate through the slides
for slide_number in range(2,11):
    slide = presentation.slides[slide_number]

    # Get the image shapes in the slide
    image_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]

    # Iterate through the image shapes
    for shape in image_shapes:
        # Get the shape name
        shape_name = shape.name

        # Get the corresponding key in the dictionary
        key = f"Slide {slide_number + 1} {shape_name}"

        # Check if the key exists in the dictionary
        if key in image_dictionary:
            # add on API token to URL
            url = image_dictionary[key] + "&apitoken=" + API_TOKEN

            response = session.get(url, verify=False)

            # save content to png file
            with open(key + ".png", "wb") as file:
                file.write(response.content)

            # Open the image
            image = Image.open(key + ".png")

            # Crop the top and bottom of the image
            cropped_image = image.crop((0, 32, image.width, image.height-10))

            # Save the cropped image as "test.png"
            cropped_image.save(key + ".png")
            file_path = key + ".png"
            
            # Remove the existing shape
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            slide.shapes._spTree.remove(shape._element)
            
            # Add the new image to the shape
            slide.shapes.add_picture(file_path, shape.left, shape.top, shape.width, shape.height)
            
            # Obtain the shape of the newly added image on the PowerPoint slide
            new_image_shape = slide.shapes[-1]

            # Change the shape name to "CPU" or "Mem" for next usage of powerpoint
            new_image_shape.name = key[-3:]

            print("Image replaced successfully. "+ key)

            #sleep for 12 seconds
            time.sleep(12)
            
# Save the modified PowerPoint presentation
presentation.save(r"file\path\modified_presentation.pptx")
